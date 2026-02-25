"""슬라이드 구조 디버깅용 스크립트"""
from pptx import Presentation
from pptx.shapes.picture import Picture

prs = Presentation('/Users/0hoxy/Downloads/★남자비모델 _수정.pptx')

for slide_idx in [1, 2, 6, 7]:
    slide = prs.slides[slide_idx - 1]
    print(f"\n{'='*60}")
    print(f"슬라이드 {slide_idx}")
    print(f"{'='*60}")

    print("\n[이미지]")
    for i, shape in enumerate(slide.shapes):
        if isinstance(shape, Picture):
            print(f"  #{i}: left={shape.left}, top={shape.top}, "
                  f"width={shape.width}, height={shape.height}, "
                  f"name='{shape.name}'")

    print("\n[텍스트]")
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame and shape.text.strip():
            print(f"  #{i}: left={shape.left}, top={shape.top}, "
                  f"text='{shape.text.strip()}'")
