"""텍스트 인코딩 디버깅"""
from pptx import Presentation

prs = Presentation('/Users/0hoxy/Downloads/★남자비모델 _수정.pptx')
slide = prs.slides[0]

for shape in slide.shapes:
    if shape.has_text_frame and shape.text.strip():
        text = shape.text.strip()
        print(f"text: '{text}'")
        print(f"repr: {repr(text)}")
        for i, ch in enumerate(text):
            print(f"  [{i}] '{ch}' -> U+{ord(ch):04X}")
