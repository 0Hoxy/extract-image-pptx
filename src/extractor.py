from pathlib import Path

from pptx import Presentation

from src.config import SLOT_NAMES
from src.classifier import filter_images, classify_slots
from src.parser import parse_slide_text
from src.storage import extract_image_data, save_image, cleanup_images


def extract_images_from_pptx(pptx_path: str, output_dir: str = "./output") -> None:
    """PPTX 파일에서 슬롯별 이미지를 추출하여 폴더에 저장한다."""
    pptx_path = Path(pptx_path)
    output_path = Path(output_dir)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX 파일을 찾을 수 없습니다: {pptx_path}")

    # PPTX 파일명으로 상위 폴더 생성
    pptx_folder = output_path / pptx_path.stem
    for slot in SLOT_NAMES:
        (pptx_folder / slot).mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    skipped = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"\n--- 슬라이드 {slide_idx} 처리 중 ---")
        if not _process_slide(slide, slide_idx, pptx_folder):
            skipped.append(slide_idx)

    print(f"\n추출 완료! 결과: {pptx_folder.resolve()}")
    print(f"총 {len(prs.slides)}장 중 {len(prs.slides) - len(skipped)}장 추출, {len(skipped)}장 제외")
    if skipped:
        print(f"제외된 슬라이드: {skipped}")


def _process_slide(slide, slide_idx: int, output_path: Path) -> bool:
    """개별 슬라이드에서 이미지를 추출하고 분류한다. 성공 시 True 반환."""
    # 이미지 필터링
    images = filter_images(slide)

    if len(images) != 5:
        print(f"  [경고] 슬라이드 {slide_idx}: 사진 {len(images)}개 발견 (5개 필요). 스킵합니다.")
        return False

    # 텍스트에서 파일명 정보 파싱
    slide_info = parse_slide_text(slide, slide_idx)
    base_name = slide_info.base_name

    # 슬롯 분류
    slot_map = classify_slots(images)

    # 이미지 데이터 준비 (all or nothing)
    prepared = []
    for slot_name, picture in slot_map.items():
        try:
            image_data = extract_image_data(picture, slot_name)
            prepared.append(image_data)
        except Exception as e:
            print(f"  [경고] 슬라이드 {slide_idx}: {slot_name} 이미지 추출 실패 ({e}). 전체 스킵합니다.")
            return False

    # 모든 이미지 저장
    saved_files = []
    for image_data in prepared:
        try:
            filepath = save_image(image_data, base_name, output_path)
            saved_files.append(filepath)
            print(f"  {image_data.slot}: {filepath.name} 저장 완료")
        except Exception as e:
            print(f"  [경고] 슬라이드 {slide_idx}: {image_data.slot} 저장 실패 ({e}). 롤백합니다.")
            cleanup_images(saved_files)
            return False

    return True
