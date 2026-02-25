import re
from pathlib import Path
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.util import Emu


SLOT_NAMES = ["MAIN", "SUB1", "SUB2", "SUB3", "SUB4"]

# 텍스트 파싱 정규식
# 형식1: "이름 (99) 178cm"
# 형식2: "VINCENT 1997's 179cm 프랑스"
TEXT_PATTERN_PAREN = re.compile(r"(.+?)\s*\((\d{2})\)\s*(\d{2,3})\s*(?:cm)?")
TEXT_PATTERN_QUOTE = re.compile(r"(.+?)\s+(\d{4})(?:['\u2019]s)?\s+(\d{2,3})\s*cm")

# 작은 이미지 필터 기준 (150pt x 200pt, 1pt = 12700 EMU)
MIN_IMAGE_WIDTH = 100 * 12700   # 1_905_000 EMU
MIN_IMAGE_HEIGHT = 100 * 12700  # 2_540_000 EMU

def extract_images_from_pptx(pptx_path: str, output_dir: str = "./output") -> None:
    """PPTX 파일에서 슬롯별 이미지를 추출하여 폴더에 저장한다."""
    pptx_path = Path(pptx_path)
    output_path = Path(output_dir)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX 파일을 찾을 수 없습니다: {pptx_path}")

    # PPTX 파일명으로 상위 폴더 생성
    pptx_folder = output_path / pptx_path.stem
    for slot in SLOT_NAMES:
        (pptx_folder / slot).mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))

    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"\n--- 슬라이드 {slide_idx} 처리 중 ---")
        _process_slide(slide, slide_idx, pptx_folder)

    print(f"\n추출 완료! 결과: {pptx_folder.resolve()}")


def _process_slide(slide, slide_idx: int, output_path: Path) -> None:
    """개별 슬라이드에서 이미지를 추출하고 분류한다."""
    # 이미지 Shape 수집 (아이콘/로고 등 작은 이미지 제외)
    images = []
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            if shape.width >= MIN_IMAGE_WIDTH and shape.height >= MIN_IMAGE_HEIGHT:
                images.append(shape)

    if len(images) != 5:
        print(f"  [경고] 슬라이드 {slide_idx}: 사진 {len(images)}개 발견 (5개 필요). 스킵합니다.")
        return

    # 텍스트에서 파일명 정보 파싱
    base_name = _parse_slide_text(slide, slide_idx)

    # 슬롯 분류
    slot_map = _classify_slots(images)

    # 이미지 데이터를 먼저 메모리에 준비 (all or nothing)
    prepared = []
    for slot_name, picture in slot_map.items():
        try:
            blob, ext = _extract_image_blob(picture)
            prepared.append((slot_name, blob, ext))
        except Exception as e:
            print(f"  [경고] 슬라이드 {slide_idx}: {slot_name} 이미지 추출 실패 ({e}). 전체 스킵합니다.")
            # 이미 저장된 파일 삭제
            for saved_slot, _, saved_ext in prepared:
                _delete_saved_image(saved_slot, base_name, saved_ext, output_path)
            return

    # 모든 이미지 준비 완료 → 저장
    for slot_name, blob, ext in prepared:
        filepath = _save_blob(blob, ext, slot_name, base_name, output_path)
        print(f"  {slot_name}: {filepath.name} 저장 완료")


def _classify_slots(images: list) -> dict:
    """이미지 5개를 좌표 기반으로 MAIN, SUB1~4로 분류한다."""
    # left 기준으로 정렬하여 가장 왼쪽 = MAIN
    sorted_by_left = sorted(images, key=lambda s: s.left)
    main_image = sorted_by_left[0]
    sub_images = sorted_by_left[1:]

    # 나머지 4개를 top 기준으로 상위/하위 분리
    sorted_by_top = sorted(sub_images, key=lambda s: s.top)
    top_row = sorted(sorted_by_top[:2], key=lambda s: s.left)
    bottom_row = sorted(sorted_by_top[2:], key=lambda s: s.left)

    return {
        "MAIN": main_image,
        "SUB1": top_row[0],
        "SUB2": top_row[1],
        "SUB3": bottom_row[0],
        "SUB4": bottom_row[1],
    }


def _parse_slide_text(slide, slide_idx: int) -> str:
    """슬라이드 하단 텍스트에서 '이름(년생)키cm' 형식의 파일명을 생성한다."""
    # 텍스트 Shape을 top 기준 내림차순 정렬 (하단부터 탐색)
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            text_shapes.append(shape)

    if not text_shapes:
        print(f"  [경고] 슬라이드 {slide_idx}: 텍스트를 찾을 수 없습니다. 기본 파일명 사용.")
        return f"slide_{slide_idx}"

    # 하단 우선, 같은 높이면 좌측 우선 (우측 특이사항보다 좌측 이름이 먼저 선택)
    text_shapes.sort(key=lambda s: (-s.top, s.left))

    for shape in text_shapes:
        text = shape.text.strip()

        # 형식1: "이름 (99) 178cm"
        match = TEXT_PATTERN_PAREN.search(text)
        if match:
            name = match.group(1).strip()
            year = match.group(2)
            height = match.group(3)
            return f"{name}({year}){height}cm"

        # 형식2: "VINCENT 1997's 179cm 프랑스"
        match = TEXT_PATTERN_QUOTE.search(text)
        if match:
            name = match.group(1).strip()
            year = match.group(2)[-2:]  # 4자리 → 2자리
            height = match.group(3)
            return f"{name}({year}){height}cm"

    # 파싱 실패 시 가장 하단 텍스트를 파일명으로 사용
    fallback = text_shapes[0].text.strip()
    # 파일명에 사용 불가한 문자 제거
    fallback = re.sub(r'[\\/:*?"<>|]', '_', fallback)
    print(f"  [경고] 슬라이드 {slide_idx}: 텍스트 파싱 실패. '{fallback}' 사용.")
    return fallback


def _extract_image_blob(picture: Picture) -> tuple[bytes, str]:
    """이미지에서 바이너리 데이터와 확장자를 추출한다."""
    image = picture.image
    content_type = image.content_type
    ext = content_type.split("/")[-1].lower()

    # 포맷 정규화
    ext_map = {"jpeg": "jpg", "x-ms-bmp": "bmp"}
    ext = ext_map.get(ext, ext)

    # MPO 등 미지원 포맷 → blob을 직접 읽어서 jpg로 저장
    supported = {"bmp", "gif", "jpg", "png", "tiff", "wmf"}
    if ext not in supported:
        ext = "jpg"

    return image.blob, ext


def _save_blob(blob: bytes, ext: str, slot_name: str, base_name: str, output_path: Path) -> Path:
    """바이너리 데이터를 슬롯 폴더에 저장한다."""
    folder = output_path / slot_name
    filename = f"{base_name}_{slot_name}.{ext}"
    filepath = folder / filename

    # 파일명 충돌 처리
    counter = 1
    while filepath.exists():
        filename = f"{base_name}({counter})_{slot_name}.{ext}"
        filepath = folder / filename
        counter += 1

    filepath.write_bytes(blob)
    return filepath


def _delete_saved_image(slot_name: str, base_name: str, ext: str, output_path: Path) -> None:
    """추출 실패 시 이미 저장된 이미지를 삭제한다."""
    folder = output_path / slot_name
    # 해당 base_name으로 시작하는 파일 찾아서 삭제
    for f in folder.glob(f"{base_name}*_{slot_name}.{ext}"):
        f.unlink()
        print(f"  [정리] {f.name} 삭제")
